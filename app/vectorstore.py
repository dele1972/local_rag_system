# app/vectorstore.py
from langchain_chroma import Chroma
from langchain_ollama.embeddings import OllamaEmbeddings
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from app.config import config
import os
import shutil
from pypdf import PdfReader
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import io
import zipfile
import xml.etree.ElementTree as ET
from openpyxl import load_workbook
import pandas as pd
from pptx import Presentation
import docx
import logging
import numpy as np
from typing import List, Dict, Any, Optional, Union

# Neue Imports für deutsche Embeddings
try:
    from langchain_huggingface import HuggingFaceEmbeddings
    HUGGINGFACE_AVAILABLE = True
except ImportError:
    HUGGINGFACE_AVAILABLE = False
    logging.warning("HuggingFace Embeddings nicht verfügbar. Installiere: pip install langchain-huggingface sentence-transformers")

try:
    from sentence_transformers import SentenceTransformer
    SENTENCE_TRANSFORMERS_AVAILABLE = True
except ImportError:
    SENTENCE_TRANSFORMERS_AVAILABLE = False
    logging.warning("SentenceTransformers nicht verfügbar. Installiere: pip install sentence-transformers")

# Logging setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Deutsche Embedding-Modelle Konfiguration
GERMAN_EMBEDDING_MODELS = {
    "huggingface": {
        "intfloat/multilingual-e5-large": {
            "model_name": "intfloat/multilingual-e5-large",
            "dimensions": 1024,
            "max_seq_length": 512,
            "german_quality": "excellent",
            "size_gb": 2.24,
            "description": "Bestes multilinguale Modell für deutsche Texte",
            "prefix": "query: ",  # Für E5 Modelle erforderlich
            "normalize_embeddings": True
        },
        "intfloat/multilingual-e5-base": {
            "model_name": "intfloat/multilingual-e5-base", 
            "dimensions": 768,
            "max_seq_length": 512,
            "german_quality": "very_good",
            "size_gb": 1.11,
            "description": "Gutes multilinguale Modell, kleiner als large",
            "prefix": "query: ",
            "normalize_embeddings": True
        },
        "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2": {
            "model_name": "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2",
            "dimensions": 384,
            "max_seq_length": 128,
            "german_quality": "good",
            "size_gb": 0.47,
            "description": "Kompaktes multilinguale Modell",
            "prefix": "",
            "normalize_embeddings": False
        },
        "sentence-transformers/paraphrase-multilingual-mpnet-base-v2": {
            "model_name": "sentence-transformers/paraphrase-multilingual-mpnet-base-v2",
            "dimensions": 768,
            "max_seq_length": 128,
            "german_quality": "very_good",
            "size_gb": 1.11,
            "description": "Sehr gutes multilinguale Modell für Paraphrasierung",
            "prefix": "",
            "normalize_embeddings": False
        }
    },
    "ollama": {
        "nomic-embed-text": {
            "model_name": "nomic-embed-text",
            "dimensions": 768,
            "max_seq_length": 8192,
            "german_quality": "fair",
            "size_gb": 0.27,
            "description": "Standard Ollama Embedding (nicht optimal für Deutsch)",
            "prefix": "",
            "normalize_embeddings": False
        },
        "mxbai-embed-large": {
            "model_name": "mxbai-embed-large",
            "dimensions": 1024,
            "max_seq_length": 512,
            "german_quality": "good",
            "size_gb": 0.67,
            "description": "Größeres Ollama Embedding mit besserer Mehrsprachigkeit",
            "prefix": "",
            "normalize_embeddings": False
        },
        "snowflake-arctic-embed": {
            "model_name": "snowflake-arctic-embed",
            "dimensions": 1024,
            "max_seq_length": 512,
            "german_quality": "good",
            "size_gb": 0.47,
            "description": "Snowflake Arctic Embedding mit guter Mehrsprachigkeit",
            "prefix": "",
            "normalize_embeddings": False
        }
    }
}

class EmbeddingManager:
    """Verwaltet verschiedene Embedding-Typen für optimale deutsche Textverarbeitung"""
    
    def __init__(self):
        self.embedding_cache = {}
        self.logger = logging.getLogger(__name__)
    
    def get_available_models(self) -> Dict[str, List[str]]:
        """Gibt verfügbare Embedding-Modelle zurück"""
        available = {"ollama": [], "huggingface": []}
        
        # Ollama Modelle (immer verfügbar)
        available["ollama"] = list(GERMAN_EMBEDDING_MODELS["ollama"].keys())
        
        # HuggingFace Modelle (nur wenn installiert)
        if HUGGINGFACE_AVAILABLE:
            available["huggingface"] = list(GERMAN_EMBEDDING_MODELS["huggingface"].keys())
        
        return available
    
    def get_model_info(self, model_type: str, model_name: str) -> Optional[Dict[str, Any]]:
        """Gibt Informationen zu einem spezifischen Modell zurück"""
        if model_type in GERMAN_EMBEDDING_MODELS and model_name in GERMAN_EMBEDDING_MODELS[model_type]:
            return GERMAN_EMBEDDING_MODELS[model_type][model_name]
        return None
    
    def get_best_german_model(self, max_size_gb: float = 3.0) -> tuple[str, str]:
        """Gibt das beste verfügbare deutsche Embedding-Modell zurück"""
        best_model = None
        best_type = None
        best_quality_score = 0
        
        quality_scores = {"excellent": 4, "very_good": 3, "good": 2, "fair": 1}
        
        for model_type, models in GERMAN_EMBEDDING_MODELS.items():
            # Überspringe HuggingFace wenn nicht verfügbar
            if model_type == "huggingface" and not HUGGINGFACE_AVAILABLE:
                continue
                
            for model_name, model_info in models.items():
                if model_info["size_gb"] <= max_size_gb:
                    quality_score = quality_scores.get(model_info["german_quality"], 0)
                    if quality_score > best_quality_score:
                        best_quality_score = quality_score
                        best_model = model_name
                        best_type = model_type
        
        # Fallback auf Ollama Standard
        if not best_model:
            best_model = "nomic-embed-text"
            best_type = "ollama"
        
        self.logger.info(f"Bestes deutsches Embedding-Modell: {best_type}/{best_model} (Qualität: {best_quality_score})")
        return best_type, best_model
    
    def create_embeddings(self, model_type: str, model_name: str) -> Union[OllamaEmbeddings, 'HuggingFaceEmbeddings']:
        """Erstellt Embedding-Objekt basierend auf Typ und Modell"""
        cache_key = f"{model_type}:{model_name}"
        
        if cache_key in self.embedding_cache:
            return self.embedding_cache[cache_key]
        
        model_info = self.get_model_info(model_type, model_name)
        if not model_info:
            raise ValueError(f"Unbekanntes Modell: {model_type}/{model_name}")
        
        try:
            if model_type == "ollama":
                embeddings = OllamaEmbeddings(
                    base_url=config.get_ollama_base_url(),
                    model=model_name
                )
                self.logger.info(f"✅ Ollama Embedding erstellt: {model_name}")
                
            elif model_type == "huggingface":
                if not HUGGINGFACE_AVAILABLE:
                    raise ImportError("HuggingFace Embeddings nicht verfügbar")
                
                # Spezielle Konfiguration für E5 Modelle
                model_kwargs = {}
                encode_kwargs = {}
                
                if "e5" in model_name.lower():
                    # E5 Modelle benötigen query/passage Präfixe
                    encode_kwargs["normalize_embeddings"] = model_info.get("normalize_embeddings", True)
                
                embeddings = HuggingFaceEmbeddings(
                    model_name=model_info["model_name"],
                    model_kwargs=model_kwargs,
                    encode_kwargs=encode_kwargs,
                    cache_folder="./models/embeddings"  # Lokaler Cache
                )
                self.logger.info(f"✅ HuggingFace Embedding erstellt: {model_name}")
                
            else:
                raise ValueError(f"Unbekannter Embedding-Typ: {model_type}")
            
            self.embedding_cache[cache_key] = embeddings
            return embeddings
            
        except Exception as e:
            self.logger.error(f"❌ Fehler beim Erstellen des Embeddings {model_type}/{model_name}: {e}")
            # Fallback auf Ollama Standard
            if model_type != "ollama" or model_name != "nomic-embed-text":
                self.logger.info("Fallback auf nomic-embed-text")
                return self.create_embeddings("ollama", "nomic-embed-text")
            raise

# Globaler Embedding Manager
embedding_manager = EmbeddingManager()

def get_embeddings(model_name: str = None, model_type: str = None):
    """
    Erstellt ein Embedding-Objekt mit optimaler deutscher Unterstützung
    
    Args:
        model_name: Spezifisches Modell (optional)
        model_type: "ollama" oder "huggingface" (optional)
    
    Returns:
        Embedding-Objekt
    """
    if model_name is None or model_type is None:
        # Automatische Auswahl des besten deutschen Modells
        model_type, model_name = embedding_manager.get_best_german_model()
        logger.info(f"🎯 Auto-Auswahl: {model_type}/{model_name}")
    
    return embedding_manager.create_embeddings(model_type, model_name)

def get_optimized_german_embeddings():
    """Gibt das optimierte deutsche Embedding-Modell zurück"""
    try:
        # Versuche zuerst das beste HuggingFace Modell
        if HUGGINGFACE_AVAILABLE:
            return embedding_manager.create_embeddings("huggingface", "intfloat/multilingual-e5-large")
        else:
            # Fallback auf bestes Ollama Modell
            return embedding_manager.create_embeddings("ollama", "mxbai-embed-large")
    except Exception as e:
        logger.warning(f"Fallback auf Standard-Embedding: {e}")
        return embedding_manager.create_embeddings("ollama", "nomic-embed-text")

def is_page_scannable(page):
    """Prüft, ob eine Seite hauptsächlich aus Bildern besteht (OCR erforderlich)"""
    text = page.extract_text()
    return not text or len(text.strip()) < 50  # Schwellenwert anpassbar

def extract_text_from_pdf_with_ocr(file_path):
    """Extrahiere Text aus PDF, mit OCR-Fallback für gescannte Seiten"""
    reader = PdfReader(file_path)
    raw_text = ""

    for i, page in enumerate(reader.pages):
        try:
            if is_page_scannable(page):
                # OCR für gescannte Seiten
                try:
                    with fitz.open(file_path) as doc:
                        pix = doc[i].get_pixmap()
                        img = Image.open(io.BytesIO(pix.tobytes("png")))
                        ocr_text = pytesseract.image_to_string(img, lang="deu")
                        if not ocr_text.strip():
                            logger.warning(f"OCR hat keinen Text auf Seite {i+1} erkannt")
                        raw_text += ocr_text + "\n"
                except pytesseract.TesseractError as e:
                    logger.error(f"OCR-Fehler auf Seite {i+1}: {e}")
                except Exception as e:
                    logger.error(f"Allgemeiner OCR-Fehler auf Seite {i+1}: {e}")
            else:
                # Normaler Text extrahierbar
                text = page.extract_text()
                raw_text += text + "\n"
        except Exception as e:
            logger.error(f"Fehler beim Verarbeiten der Seite {i+1}: {e}")
    
    return raw_text

def extract_text_from_docx(file_path):
    """Extrahiere Text aus Word-Dokument (.docx)"""
    try:
        doc = docx.Document(file_path)
        text = ""
        
        # Haupttext
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Text aus Tabellen
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + "\t"
                text += "\n"
        
        # Text aus Kopf- und Fußzeilen
        for section in doc.sections:
            if section.header:
                for paragraph in section.header.paragraphs:
                    text += paragraph.text + "\n"
            if section.footer:
                for paragraph in section.footer.paragraphs:
                    text += paragraph.text + "\n"
                    
        return text
    except Exception as e:
        logger.error(f"Fehler beim Extrahieren aus Word-Dokument {file_path}: {e}")
        return ""

def extract_text_from_doc(file_path):
    """Extrahiere Text aus altem Word-Format (.doc) über LibreOffice"""
    try:
        # Fallback für .doc Dateien - erfordert LibreOffice
        import subprocess
        import tempfile
        
        with tempfile.TemporaryDirectory() as temp_dir:
            # Konvertiere .doc zu .txt mit LibreOffice
            output_file = os.path.join(temp_dir, "converted.txt")
            subprocess.run([
                "libreoffice", "--headless", "--convert-to", "txt",
                "--outdir", temp_dir, file_path
            ], check=True, capture_output=True)
            
            # Lese konvertierte Datei
            txt_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(file_path))[0] + ".txt")
            if os.path.exists(txt_file):
                with open(txt_file, 'r', encoding='utf-8') as f:
                    return f.read()
            return ""
    except Exception as e:
        logger.error(f"Fehler beim Extrahieren aus .doc Datei {file_path}: {e}")
        logger.info("Hinweis: .doc Dateien benötigen LibreOffice für die Konvertierung")
        return ""

def extract_text_from_excel(file_path):
    """Extrahiere Text aus Excel-Datei (.xlsx, .xls)"""
    try:
        text = ""
        
        # Verwende pandas für bessere Kompatibilität
        if file_path.endswith('.xlsx'):
            # Lade alle Sheets
            xl_file = pd.ExcelFile(file_path)
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text += f"\n=== Sheet: {sheet_name} ===\n"
                
                # Konvertiere DataFrame zu Text
                for index, row in df.iterrows():
                    row_text = []
                    for col in df.columns:
                        cell_value = str(row[col]) if pd.notna(row[col]) else ""
                        if cell_value:
                            row_text.append(f"{col}: {cell_value}")
                    if row_text:
                        text += " | ".join(row_text) + "\n"
        
        elif file_path.endswith('.xls'):
            # Für .xls Dateien
            df = pd.read_excel(file_path, engine='xlrd')
            for index, row in df.iterrows():
                row_text = []
                for col in df.columns:
                    cell_value = str(row[col]) if pd.notna(row[col]) else ""
                    if cell_value:
                        row_text.append(f"{col}: {cell_value}")
                if row_text:
                    text += " | ".join(row_text) + "\n"
        
        return text
    except Exception as e:
        logger.error(f"Fehler beim Extrahieren aus Excel-Datei {file_path}: {e}")
        return ""

def extract_text_from_powerpoint(file_path):
    """Extrahiere Text aus PowerPoint-Präsentation (.pptx)"""
    try:
        prs = Presentation(file_path)
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n=== Folie {slide_num} ===\n"
            
            # Text aus Shapes extrahieren
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
                
                # Text aus Tabellen in Shapes
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text += " | ".join(row_text) + "\n"
            
            # Notizen extrahieren
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                    text += f"Notizen: {notes_slide.notes_text_frame.text}\n"
        
        return text
    except Exception as e:
        logger.error(f"Fehler beim Extrahieren aus PowerPoint-Datei {file_path}: {e}")
        return ""

def extract_text_from_ppt(file_path):
    """Extrahiere Text aus altem PowerPoint-Format (.ppt) über LibreOffice"""
    try:
        import subprocess
        import tempfile
        
        with tempfile.TemporaryDirectory() as temp_dir:
            # Konvertiere .ppt zu .txt mit LibreOffice
            subprocess.run([
                "libreoffice", "--headless", "--convert-to", "txt",
                "--outdir", temp_dir, file_path
            ], check=True, capture_output=True)
            
            # Lese konvertierte Datei
            txt_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(file_path))[0] + ".txt")
            if os.path.exists(txt_file):
                with open(txt_file, 'r', encoding='utf-8') as f:
                    return f.read()
            return ""
    except Exception as e:
        logger.error(f"Fehler beim Extrahieren aus .ppt Datei {file_path}: {e}")
        logger.info("Hinweis: .ppt Dateien benötigen LibreOffice für die Konvertierung")
        return ""

def get_optimized_chunk_settings(embedding_model_info: Dict[str, Any], file_size_mb: float = None) -> Dict[str, int]:
    """
    Bestimmt optimale Chunk-Einstellungen basierend auf Embedding-Modell und Dateigröße
    
    Args:
        embedding_model_info: Informationen zum Embedding-Modell
        file_size_mb: Dateigröße in MB
    
    Returns:
        Dict mit chunk_size und chunk_overlap
    """
    max_seq_length = embedding_model_info.get("max_seq_length", 512)
    
    # Grundeinstellungen basierend auf Modell-Kapazität
    if max_seq_length >= 8192:
        # Große Modelle können längere Chunks verarbeiten
        base_chunk_size = 2000
        base_overlap = 400
    elif max_seq_length >= 512:
        # Standard Modelle
        base_chunk_size = 1000
        base_overlap = 200
    else:
        # Kleine Modelle
        base_chunk_size = 500
        base_overlap = 100
    
    # Anpassung basierend auf Dateigröße
    if file_size_mb:
        if file_size_mb > 5.0:
            # Große Dateien: Kleinere Chunks für bessere Präzision
            chunk_size = int(base_chunk_size * 0.8)
            overlap = int(base_overlap * 1.2)
        elif file_size_mb < 0.5:
            # Kleine Dateien: Größere Chunks für besseren Kontext
            chunk_size = int(base_chunk_size * 1.2)
            overlap = int(base_overlap * 0.8)
        else:
            chunk_size = base_chunk_size
            overlap = base_overlap
    else:
        chunk_size = base_chunk_size
        overlap = base_overlap
    
    return {
        "chunk_size": chunk_size,
        "chunk_overlap": overlap
    }

def chunk_text_german_optimized(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    """
    Deutsche-optimierte Text-Chunking-Funktion
    
    Args:
        text: Zu chunkender Text
        chunk_size: Gewünschte Chunk-Größe
        chunk_overlap: Überlappung zwischen Chunks
    
    Returns:
        Liste von Text-Chunks
    """
    # Deutsche Separatoren für bessere Textaufteilung
    german_separators = [
        "\n\n",  # Absätze
        "\n",    # Zeilenwechsel
        ". ",    # Satzende
        "? ",    # Fragen
        "! ",    # Ausrufe
        "; ",    # Semikolon
        ", ",    # Komma
        " ",     # Leerzeichen
        ""       # Fallback
    ]
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=german_separators,
        keep_separator=True  # Wichtig für deutschen Text
    )
    
    chunks = text_splitter.split_text(text)
    
    # Post-Processing: Entferne zu kleine oder leere Chunks
    filtered_chunks = []
    min_chunk_size = max(50, int(chunk_size * 0.1))  # Mindestens 50 Zeichen oder 10% der Chunk-Größe
    
    for chunk in chunks:
        cleaned_chunk = chunk.strip()
        if len(cleaned_chunk) >= min_chunk_size:
            filtered_chunks.append(cleaned_chunk)
    
    logger.info(f"📊 Text-Chunking: {len(chunks)} → {len(filtered_chunks)} Chunks (nach Filterung)")
    return filtered_chunks

def chunk_text(text, chunk_size=1000, chunk_overlap=200):
    """Legacy-Funktion für Rückwärtskompatibilität"""
    return chunk_text_german_optimized(text, chunk_size, chunk_overlap)

def load_documents_from_file(file_path, chunk_size=1000, chunk_overlap=200):
    """Lädt Text aus verschiedenen Dateiformaten mit optimiertem Chunking"""
    content = ""
    file_extension = os.path.splitext(file_path)[1].lower()
    
    logger.info(f"📄 Verarbeitung: {file_path}")
    try:
        if file_extension in [".txt", ".md"]:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
        
        elif file_extension == ".pdf":
            content = extract_text_from_pdf_with_ocr(file_path)
        
        elif file_extension == ".docx":
            content = extract_text_from_docx(file_path)
        
        elif file_extension == ".doc":
            content = extract_text_from_doc(file_path)
        
        elif file_extension in [".xlsx", ".xls"]:
            content = extract_text_from_excel(file_path)
        
        elif file_extension == ".pptx":
            content = extract_text_from_powerpoint(file_path)
        
        elif file_extension == ".ppt":
            content = extract_text_from_ppt(file_path)
        
        else:
            raise ValueError(f"Nicht unterstützter Dateityp: {file_extension}")
        
        if not content or len(content.strip()) == 0:
            logger.warning(f"⚠️ Keine Inhalte aus Datei {file_path} extrahiert")
            return []
        
        # Dateigröße für Chunk-Optimierung berechnen
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        
        # Optimierte Chunk-Einstellungen
        if len(content) > chunk_size:
            chunks = chunk_text_german_optimized(content, chunk_size, chunk_overlap)
            documents = []
            for i, chunk in enumerate(chunks):
                if chunk.strip():  # Nur nicht-leere Chunks
                    doc = Document(
                        page_content=chunk,
                        metadata={
                            "source": file_path,
                            "file_type": file_extension,
                            "chunk": i,
                            "total_chunks": len(chunks),
                            "file_size_mb": round(file_size_mb, 2),
                            "chunk_method": "german_optimized"
                        }
                    )
                    documents.append(doc)
            
            logger.info(f"✅ {len(documents)} Chunks aus {file_path} erstellt ({file_size_mb:.1f}MB)")
            return documents
        else:
            return [Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "file_type": file_extension,
                    "file_size_mb": round(file_size_mb, 2),
                    "chunk_method": "single_document"
                }
            )]
    
    except Exception as e:
        logger.error(f"❌ Fehler beim Laden der Datei {file_path}: {e}")
        return []

def load_documents(documents_path, chunk_size=1000, chunk_overlap=200):
    """
    Lädt alle Dokumente aus einem Pfad (Datei oder Verzeichnis) mit deutscher Optimierung
    
    Args:
        documents_path: Pfad zu einer Datei oder einem Verzeichnis
        chunk_size: Größe der Text-Chunks  
        chunk_overlap: Überlappung zwischen Chunks
    
    Returns:
        Liste von LangChain-Dokumenten
    """
    documents = []
    supported_extensions = config.get_supported_file_types()
    
    if os.path.isfile(documents_path):
        # Einzelne Datei
        logger.info(f"📁 Lade Datei: {documents_path}")
        documents.extend(load_documents_from_file(documents_path, chunk_size, chunk_overlap))
        
    elif os.path.isdir(documents_path):
        # Verzeichnis durchsuchen
        logger.info(f"📁 Durchsuche Verzeichnis: {documents_path}")
        file_count = 0
        total_size_mb = 0
        
        for root, dirs, files in os.walk(documents_path):
            for file in files:
                file_path = os.path.join(root, file)
                file_extension = os.path.splitext(file)[1].lower()
                
                if file_extension in supported_extensions:
                    file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
                    total_size_mb += file_size_mb
                    
                    logger.info(f"📄 Lade: {file_path} ({file_size_mb:.1f}MB)")
                    file_documents = load_documents_from_file(file_path, chunk_size, chunk_overlap)
                    documents.extend(file_documents)
                    file_count += 1
                else:
                    logger.debug(f"⏭️ Überspringe: {file_path} (nicht unterstützt)")
        
        logger.info(f"📊 Zusammenfassung: {file_count} Dateien, {total_size_mb:.1f}MB, {len(documents)} Chunks")
    else:
        raise ValueError(f"Pfad nicht gefunden: {documents_path}")
    
    if not documents:
        logger.warning("⚠️ Keine Dokumente gefunden oder alle Dateien waren leer")
    
    return documents

def analyze_similarity_scores(vectorstore, query: str, k: int = 10) -> Dict[str, Any]:
    """
    Analysiert Similarity-Scores für bessere Retrieval-Optimierung
    
    Args:
        vectorstore: Chroma-Vectorstore
        query: Suchanfrage
        k: Anzahl der zu retrievenden Dokumente
    
    Returns:
        Analyse-Ergebnisse mit Scores und Empfehlungen
    """
    try:
        # Hole Dokumente mit Scores
        results = vectorstore.similarity_search_with_score(query, k=k)
        
        if not results:
            return {
                "scores": [],
                "recommendations": {
                    "threshold": 0.7,
                    "optimal_k": 5,
                    "quality": "no_results"
                }
            }
        
        scores = [score for _, score in results]
        
        # Statistiken berechnen
        min_score = min(scores)
        max_score = max(scores)
        avg_score = np.mean(scores)
        std_score = np.std(scores)
        
        # Qualitätsbewertung
        if avg_score < 0.5:
            quality = "excellent"
            recommended_threshold = 0.7
        elif avg_score < 0.8:
            quality = "good"
            recommended_threshold = 0.8
        elif avg_score < 1.2:
            quality = "fair"
            recommended_threshold = 1.0
        else:
            quality = "poor"
            recommended_threshold = 1.5
        
        # Optimale K-Wert Bestimmung (wo Scores stark abfallen)
        optimal_k = k
        if len(scores) > 3:
            score_diffs = [scores[i+1] - scores[i] for i in range(len(scores)-1)]
            if score_diffs:
                max_diff_idx = score_diffs.index(max(score_diffs))
                optimal_k = min(max_diff_idx + 2, k)  # +2 weil wir nach dem größten Sprung schauen
        
        analysis_result = {
            "scores": scores,
            "statistics": {
                "min": min_score,
                "max": max_score,
                "avg": avg_score,
                "std": std_score,
                "count": len(scores)
            },
            "recommendations": {
                "threshold": recommended_threshold,
                "optimal_k": optimal_k,
                "quality": quality
            },
            "documents": [
                {
                    "content_preview": doc.page_content[:100] + "...",
                    "score": score,
                    "source": doc.metadata.get("source", "unknown"),
                    "chunk": doc.metadata.get("chunk", 0)
                }
                for doc, score in results[:5]  # Top 5 für Analyse
            ]
        }
        
        logger.info(f"🎯 Similarity-Analyse: {quality} (Ø {avg_score:.3f}, optimal k={optimal_k})")
        return analysis_result
        
    except Exception as e:
        logger.error(f"❌ Fehler bei Similarity-Analyse: {e}")
        return {
            "scores": [],
            "recommendations": {
                "threshold": 0.8,
                "optimal_k": 5,
                "quality": "error"
            }
        }

def build_vectorstore_with_analysis(documents, model_name=None, model_type=None, persist_directory=None):
    """Erstelle einen Vektorspeicher mit deutscher Optimierung und Analyse"""
    if not documents:
        raise ValueError("Keine Dokumente zum Erstellen des Vektorspeichers vorhanden")
    
    # Automatische Modell-Auswahl wenn nicht spezifiziert
    if model_name is None or model_type is None:
        model_type, model_name = embedding_manager.get_best_german_model()
        logger.info(f"🎯 Auto-Auswahl Embedding: {model_type}/{model_name}")
    
    # Embedding-Modell erstellen
    embeddings = embedding_manager.create_embeddings(model_type, model_name)
    
    # Modell-Informationen für Optimierung
    model_info = embedding_manager.get_model_info(model_type, model_name)
    if model_info:
        logger.info(f"📊 Modell-Info: {model_info['dimensions']}D, {model_info['max_seq_length']} tokens, {model_info['german_quality']} Deutsch-Qualität")
    
    # Vektorspeicher erstellen
    try:
        if persist_directory:
            vectorstore = Chroma.from_documents(
                documents=documents, 
                embedding=embeddings,
                persist_directory=persist_directory
            )
            logger.info(f"💾 Vektorspeicher gespeichert: {persist_directory}")
        else:
            vectorstore = Chroma.from_documents(
                documents=documents, 
                embedding=embeddings
            )
            logger.info("🔄 Temporärer Vektorspeicher erstellt")
        
        # Qualitätsanalyse mit Testabfrage
        if len(documents) > 0:
            # Verwende ersten Chunk als Testabfrage
            test_query = documents[0].page_content[:100]
            analysis = analyze_similarity_scores(vectorstore, test_query, k=min(10, len(documents)))
            
            logger.info(f"📈 Vektorspeicher-Qualität: {analysis['recommendations']['quality']}")
            logger.info(f"📊 Empfohlener Threshold: {analysis['recommendations']['threshold']}")
            logger.info(f"🎯 Optimaler k-Wert: {analysis['recommendations']['optimal_k']}")
            
            # Speichere Analyse-Ergebnisse in Metadata
            if hasattr(vectorstore, '_collection'):
                vectorstore._analysis_results = analysis
        
        return vectorstore
        
    except Exception as e:
        logger.error(f"❌ Fehler beim Erstellen des Vektorspeichers: {e}")
        raise

def build_vectorstore(documents, model_name, persist_directory=None):
    """Legacy-Funktion für Rückwärtskompatibilität"""
    # Extrahiere Modell-Typ aus model_name falls möglich
    if "/" in model_name:
        # HuggingFace Format
        return build_vectorstore_with_analysis(documents, model_name, "huggingface", persist_directory)
    else:
        # Ollama Format
        return build_vectorstore_with_analysis(documents, model_name, "ollama", persist_directory)

def load_vectorstore_with_analysis(persist_directory, model_name=None, model_type=None):
    """Lade einen gespeicherten Vektorspeicher mit Analyse-Capabilities"""
    if not os.path.exists(persist_directory):
        raise FileNotFoundError(f"Vector store not found at {persist_directory}")
    
    # Automatische Modell-Auswahl wenn nicht spezifiziert
    if model_name is None or model_type is None:
        model_type, model_name = embedding_manager.get_best_german_model()
        logger.info(f"🎯 Auto-Auswahl beim Laden: {model_type}/{model_name}")
    
    embeddings = embedding_manager.create_embeddings(model_type, model_name)
    
    vectorstore = Chroma(
        persist_directory=persist_directory,
        embedding_function=embeddings
    )
    
    logger.info(f"📂 Vektorspeicher geladen: {persist_directory}")
    return vectorstore

def load_vectorstore(persist_directory, model_name):
    """Legacy-Funktion für Rückwärtskompatibilität"""
    # Extrahiere Modell-Typ aus model_name falls möglich
    if "/" in model_name:
        # HuggingFace Format
        return load_vectorstore_with_analysis(persist_directory, model_name, "huggingface")
    else:
        # Ollama Format
        return load_vectorstore_with_analysis(persist_directory, model_name, "ollama")

def delete_vectorstore(persist_directory):
    """Lösche einen Vektorspeicher"""
    if os.path.exists(persist_directory):
        shutil.rmtree(persist_directory)
        logger.info(f"🗑️ Vektorspeicher gelöscht: {persist_directory}")
        return True
    return False

def list_vectorstores(base_directory):
    """Liste alle verfügbaren Vektorspeicher in einem Verzeichnis auf"""
    vectorstores = {}
    
    if not os.path.exists(base_directory):
        return vectorstores
    
    for item in os.listdir(base_directory):
        item_path = os.path.join(base_directory, item)
        if os.path.isdir(item_path):
            chroma_files = ["chroma.sqlite3", "chroma_metadata.parquet"]
            has_chroma_files = any(os.path.exists(os.path.join(item_path, cf)) for cf in chroma_files)
            if has_chroma_files:
                vectorstores[item] = item_path
    
    return vectorstores

def get_vectorstore_info(persist_directory):
    """Gibt Informationen über einen Vektorspeicher zurück"""
    if not os.path.exists(persist_directory):
        return None
    
    try:
        # Lade temporär für Informationen
        temp_embeddings = get_optimized_german_embeddings()
        vectorstore = Chroma(
            persist_directory=persist_directory,
            embedding_function=temp_embeddings
        )
        
        # Sammle Informationen
        collection = vectorstore._collection
        count = collection.count()
        
        # Beispiel-Dokument für Metadata-Analyse
        if count > 0:
            sample_results = vectorstore.similarity_search("test", k=1)
            sample_metadata = sample_results[0].metadata if sample_results else {}
        else:
            sample_metadata = {}
        
        info = {
            "document_count": count,
            "persist_directory": persist_directory,
            "sample_metadata": sample_metadata,
            "supported_file_types": list(set(
                sample_metadata.get("file_type", "unknown") 
                for doc in vectorstore.similarity_search("", k=min(10, count))
            )) if count > 0 else [],
            "has_analysis": hasattr(vectorstore, '_analysis_results')
        }
        
        return info
        
    except Exception as e:
        logger.error(f"❌ Fehler beim Abrufen von Vektorspeicher-Informationen: {e}")
        return {"error": str(e)}

def optimize_retrieval_for_large_documents(vectorstore, query: str, file_size_mb: float = None, max_tokens: int = 4096) -> Dict[str, Any]:
    """
    Optimiert Retrieval-Parameter für große deutsche Dokumente
    
    Args:
        vectorstore: Chroma-Vectorstore
        query: Suchanfrage
        file_size_mb: Dateigröße in MB
        max_tokens: Maximale Token-Anzahl für das Modell
    
    Returns:
        Optimierte Retrieval-Parameter und Ergebnisse
    """
    logger.info(f"👉 NEU!!! optimize_retrieval_for_large_documents")
    # Basis-Parameter
    base_k = 5
    base_threshold = 0.8
    
    # Anpassung basierend auf Dateigröße
    if file_size_mb:
        if file_size_mb > 5.0:
            # Große Dateien: Mehr Dokumente, niedrigerer Threshold
            k = min(15, int(base_k * 3))
            threshold = base_threshold * 0.9
        elif file_size_mb > 2.0:
            # Mittlere Dateien
            k = min(10, int(base_k * 2))
            threshold = base_threshold * 0.95
        else:
            # Kleine Dateien
            k = base_k
            threshold = base_threshold
    else:
        k = base_k
        threshold = base_threshold
    
    # Similarity-Analyse für weitere Optimierung
    analysis = analyze_similarity_scores(vectorstore, query, k=k*2)  # Analysiere mehr als wir brauchen
    
    if analysis["recommendations"]["quality"] in ["excellent", "good"]:
        # Gute Qualität: Verwende empfohlene Werte
        optimal_k = min(analysis["recommendations"]["optimal_k"], k)
        optimal_threshold = analysis["recommendations"]["threshold"]
    else:
        # Schlechte Qualität: Konservative Werte
        optimal_k = max(3, k // 2)
        optimal_threshold = threshold * 0.8
    
    # Token-Budget berücksichtigen
    estimated_tokens_per_doc = 300  # Durchschnittliche Schätzung
    max_docs_by_tokens = max_tokens // estimated_tokens_per_doc
    final_k = min(optimal_k, max_docs_by_tokens)
    
    # Retrieval durchführen
    try:
        results = vectorstore.similarity_search_with_score(query, k=final_k)
        
        # Filtere nach Threshold
        filtered_results = [
            (doc, score) for doc, score in results 
            if score <= optimal_threshold
        ]
        
        if not filtered_results and results:
            # Fallback: Nimm beste Ergebnisse auch wenn über Threshold
            filtered_results = results[:max(1, final_k // 2)]
            logger.warning(f"⚠️ Threshold-Filter zu streng, verwende {len(filtered_results)} beste Ergebnisse")
        
        optimization_info = {
            "parameters": {
                "k": final_k,
                "threshold": optimal_threshold,
                "file_size_mb": file_size_mb,
                "max_tokens": max_tokens
            },
            "analysis": analysis,
            "results_count": len(filtered_results),
            "optimization_applied": True
        }
        
        logger.info(f"🎯 Retrieval optimiert: k={final_k}, threshold={optimal_threshold:.3f}, {len(filtered_results)} Ergebnisse")
        
        return {
            "documents": [doc for doc, _ in filtered_results],
            "scores": [score for _, score in filtered_results],
            "optimization_info": optimization_info
        }
        
    except Exception as e:
        logger.error(f"❌ Fehler bei optimiertem Retrieval: {e}")
        # Fallback auf einfaches Retrieval
        fallback_docs = vectorstore.similarity_search(query, k=3)
        return {
            "documents": fallback_docs,
            "scores": [],
            "optimization_info": {"optimization_applied": False, "error": str(e)}
        }

def get_vectorstore(documents, model_name="auto", persist_directory=None, chunk_size=1000, chunk_overlap=200):
    """
    Erstellt oder lädt einen optimierten deutschen Vektorspeicher
    
    Args:
        documents: Liste von LangChain-Dokumenten oder Pfad zu Dokumenten
        model_name: Name des Embedding-Modells ("auto" für automatische Auswahl)
        persist_directory: Verzeichnis zum Speichern des Vektorspeichers
        chunk_size: Größe der Text-Chunks
        chunk_overlap: Überlappung zwischen Chunks
    
    Returns:
        Optimierter Chroma-Vektorspeicher
    """
    
    # Falls documents ein String ist (Pfad), lade Dokumente aus dem Pfad
    if isinstance(documents, str):
        documents_path = documents
        documents = load_documents(documents_path, chunk_size, chunk_overlap)
    
    if not documents:
        raise ValueError("Keine Dokumente zum Erstellen des Vektorspeichers gefunden")
    
    # Automatische Modell-Auswahl
    if model_name == "auto":
        model_type, model_name = embedding_manager.get_best_german_model()
        logger.info(f"🎯 Automatische Modell-Auswahl: {model_type}/{model_name}")
    else:
        # Legacy-Support: Bestimme Typ aus Name
        if "/" in model_name:
            model_type = "huggingface"
        else:
            model_type = "ollama"
    
    logger.info(f"📊 Erstelle Vektorspeicher mit {len(documents)} Dokumenten")
    
    # Vektorspeicher mit Analyse erstellen
    return build_vectorstore_with_analysis(documents, model_name, model_type, persist_directory)

# Convenience-Funktionen für einfache Nutzung

def create_german_vectorstore(documents_path: str, persist_directory: str = None) -> Chroma:
    """
    Einfache Funktion zum Erstellen eines deutschen Vektorspeichers
    
    Args:
        documents_path: Pfad zu Dokumenten
        persist_directory: Speicherort (optional)
    
    Returns:
        Optimierter Vektorspeicher
    """
    logger.info(f"👉 NEU!!! create_german_vectorstore")
    return get_vectorstore(documents_path, model_name="auto", persist_directory=persist_directory)

def load_german_vectorstore(persist_directory: str) -> Chroma:
    """
    Einfache Funktion zum Laden eines deutschen Vektorspeichers
    
    Args:
        persist_directory: Speicherort
    
    Returns:
        Geladener Vektorspeicher
    """
    logger.info(f"👉 NEU!!! load_german_vectorstore")
    return load_vectorstore_with_analysis(persist_directory)

def get_embedding_model_recommendations() -> Dict[str, Any]:
    """Gibt Empfehlungen für deutsche Embedding-Modelle zurück"""
    available = embedding_manager.get_available_models()
    
    recommendations = {
        "best_overall": None,
        "best_quality": None, 
        "best_speed": None,
        "available_models": available,
        "installation_notes": []
    }
    
    if available["huggingface"]:
        recommendations["best_overall"] = ("huggingface", "intfloat/multilingual-e5-large")
        recommendations["best_quality"] = ("huggingface", "intfloat/multilingual-e5-large")
        recommendations["best_speed"] = ("huggingface", "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2")
    else:
        recommendations["best_overall"] = ("ollama", "mxbai-embed-large")
        recommendations["best_quality"] = ("ollama", "mxbai-embed-large") 
        recommendations["best_speed"] = ("ollama", "nomic-embed-text")
        recommendations["installation_notes"].append(
            "Für bessere deutsche Unterstützung installiere: pip install langchain-huggingface sentence-transformers"
        )
    
    return recommendations